import os
import re
import csv
from pathlib import Path
from typing import List, Optional

import fitz
import pandas as pd
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation


SUPPORTED_EXTENSIONS = {
    ".pdf": "pdf",
    ".csv": "csv",
    ".docx": "docx",
    ".txt": "txt",
    ".xls": "xls",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
}


def clean_extracted_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\r\n", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+\n", "\n", text)
    return text.strip()


def extract_pdf_text(file_path: str) -> str:
    doc = fitz.open(file_path)
    chunks = []
    for page in doc:
        text = page.get_text("text")
        if text:
            chunks.append(text)
    doc.close()
    return clean_extracted_text("\n\n".join(chunks))


def extract_docx_text(file_path: str) -> str:
    document = DocxDocument(file_path)
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    return clean_extracted_text("\n".join(paragraphs))


def extract_txt_text(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as file:
        return clean_extracted_text(file.read())


def extract_csv_text(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8-sig", errors="replace", newline="") as file:
        rows = csv.reader(file)
        return clean_extracted_text("\n".join(" | ".join(cell.strip() for cell in row) for row in rows))


def extract_xlsx_text(file_path: str) -> str:
    workbook = load_workbook(file_path, read_only=True, data_only=True)
    sheet_text = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            values = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if values:
                sheet_text.append(" | ".join(values))
    workbook.close()
    return clean_extracted_text("\n".join(sheet_text))


def extract_xls_text(file_path: str) -> str:
    dataframe = pd.read_excel(file_path, engine="xlrd")
    dataframe = dataframe.fillna("")
    rows = [" | ".join(str(value).strip() for value in row) for row in dataframe.astype(str).values]
    headers = " | ".join(str(column).strip() for column in dataframe.columns)
    return clean_extracted_text("\n".join([headers, *rows]))


def extract_pptx_text(file_path: str) -> str:
    presentation = Presentation(file_path)
    slide_text = []
    for slide in presentation.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        if texts:
            slide_text.append("\n".join(texts))
    return clean_extracted_text("\n\n".join(slide_text))


def extract_text_by_file_type(file_path: str, file_type: str) -> str:
    type_map = {
        "csv": extract_csv_text,
        "pdf": extract_pdf_text,
        "docx": extract_docx_text,
        "txt": extract_txt_text,
        "xlsx": extract_xlsx_text,
        "xls": extract_xls_text,
        "pptx": extract_pptx_text,
    }
    handler = type_map.get(file_type.lower().lstrip("."))
    if handler:
        return handler(file_path)

    # Unknown extensions are still accepted and can be indexed when they contain text.
    try:
        with open(file_path, "r", encoding="utf-8-sig", errors="strict") as file:
            return clean_extracted_text(file.read())
    except (UnicodeDecodeError, OSError) as exc:
        raise ValueError(f"No text extractor is available for file type: {file_type}") from exc


def extract_structured_data(file_path: str, file_type: str) -> list[dict]:
    """Return real rows for tabular uploads; other document types have no row dataset."""
    if file_type == "csv":
        with open(file_path, "r", encoding="utf-8-sig", errors="replace", newline="") as file:
            return list(csv.DictReader(file))
    if file_type in {"xls", "xlsx"}:
        if file_type == "xls":
            dataframe = pd.read_excel(file_path, engine="xlrd")
            return dataframe.fillna("").to_dict(orient="records")
        workbook = load_workbook(file_path, read_only=True, data_only=True)
        rows = []
        for sheet in workbook.worksheets:
            values = list(sheet.iter_rows(values_only=True))
            if not values:
                continue
            headers = [str(value).strip() if value is not None else f"column_{index + 1}" for index, value in enumerate(values[0])]
            for row in values[1:]:
                rows.append({header: value for header, value in zip(headers, row)})
        workbook.close()
        return rows
    return []


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 120) -> List[str]:
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunk = text[start:end]
        chunks.append(chunk)
        if end >= len(text):
            break
        start = max(0, end - overlap)
    return chunks
